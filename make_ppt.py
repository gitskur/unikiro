from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W, H = Inches(13.33), Inches(7.5)  # 와이드 16:9

DARK_BLUE  = RGBColor(0x1F, 0x39, 0x64)
MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
ORANGE     = RGBColor(0xED, 0x7D, 0x31)
GREEN      = RGBColor(0x70, 0xAD, 0x47)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY  = RGBColor(0x40, 0x40, 0x40)


def add_text(tf, text, size, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_box(slide, x, y, w, h, text, font_size=11, bg=LIGHT_BLUE,
            fg=DARK_BLUE, bold=False, align=PP_ALIGN.CENTER, radius=True):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
    if radius:
        shape = slide.shapes.add_shape(5, x, y, w, h)  # 5 = rounded rectangle
    else:
        shape = slide.shapes.add_shape(1, x, y, w, h)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = MID_BLUE
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    add_text(tf, text, font_size, bold=bold, color=fg, align=align)
    return shape


def arrow(slide, x1, y1, x2, y2):
    """단순 선 화살표."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1=straight
    connector.line.color.rgb = MID_BLUE
    connector.line.width = Pt(1.5)


prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # 완전 빈 레이아웃


# ── 슬라이드 1: 표지 ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg = s.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BLUE

tb = s.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.2))
add_text(tb.text_frame, "Coupa × Teams 연동 챗봇", 40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tb2 = s.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10), Inches(0.7))
add_text(tb2.text_frame, "구매·비용처리 프로세스 자동화 아키텍처", 22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

tb3 = s.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.5))
add_text(tb3.text_frame, "AWS Lambda  ·  Amazon Bedrock  ·  Microsoft Teams  ·  Coupa", 14,
         color=RGBColor(0xA0, 0xC4, 0xE8), align=PP_ALIGN.CENTER)


# ── 슬라이드 2: 전체 시스템 구성 ─────────────────────────────────
s = prs.slides.add_slide(blank)
s.background.fill.solid()
s.background.fill.fore_color.rgb = GRAY

# 제목
tb = s.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12), Inches(0.55))
add_text(tb.text_frame, "전체 시스템 구성", 24, bold=True, color=DARK_BLUE)

# 영역 박스 3개
def zone(slide, x, y, w, h, label, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # 1=rectangle
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.color.rgb = MID_BLUE; shape.line.width = Pt(1)
    shape.fill.fore_color.theme_color  # no-op, just access
    tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.35))
    add_text(tb.text_frame, label, 11, bold=True, color=DARK_BLUE)

zone(s, Inches(0.3),  Inches(0.85), Inches(2.8), Inches(5.8), "💬 Microsoft Teams", RGBColor(0xE8,0xF0,0xFE))
zone(s, Inches(3.4),  Inches(0.85), Inches(6.5), Inches(5.8), "☁️ AWS",             RGBColor(0xFF,0xF3,0xE0))
zone(s, Inches(10.2), Inches(0.85), Inches(2.9), Inches(5.8), "🛒 Coupa",           RGBColor(0xE8,0xF5,0xE9))

# Teams 내부
add_box(s, Inches(0.5),  Inches(1.5),  Inches(2.3), Inches(0.6), "사용자",          12, bg=MID_BLUE,   fg=WHITE, bold=True)
add_box(s, Inches(0.5),  Inches(3.0),  Inches(2.3), Inches(0.6), "팀 채널",         12, bg=MID_BLUE,   fg=WHITE, bold=True)

# AWS 내부
add_box(s, Inches(3.6),  Inches(1.4),  Inches(2.8), Inches(0.7), "Lambda 1\nTeams Bot Handler",      10, bg=ORANGE, fg=WHITE, bold=True)
add_box(s, Inches(3.6),  Inches(2.7),  Inches(2.8), Inches(0.7), "Lambda 2\nContract Scheduler",     10, bg=ORANGE, fg=WHITE, bold=True)
add_box(s, Inches(3.6),  Inches(4.0),  Inches(2.8), Inches(0.7), "Lambda 3\nWebhook Handler",        10, bg=ORANGE, fg=WHITE, bold=True)
add_box(s, Inches(6.8),  Inches(1.4),  Inches(2.8), Inches(0.7), "Amazon Bedrock\nClaude 3.5 Sonnet",10, bg=GREEN,  fg=WHITE, bold=True)
add_box(s, Inches(6.8),  Inches(2.7),  Inches(2.8), Inches(0.7), "EventBridge\n매일 00:00 UTC",      10, bg=GREEN,  fg=WHITE, bold=True)
add_box(s, Inches(6.8),  Inches(4.0),  Inches(2.8), Inches(0.7), "DynamoDB\nChannelMapping",         10, bg=GREEN,  fg=WHITE, bold=True)

# Coupa 내부
add_box(s, Inches(10.4), Inches(1.4),  Inches(2.4), Inches(0.7), "Coupa REST API\n(/api/contracts)",  10, bg=RGBColor(0x21,0x96,0x53), fg=WHITE, bold=True)
add_box(s, Inches(10.4), Inches(2.7),  Inches(2.4), Inches(0.7), "Coupa Webhook\n(상태변경 이벤트)", 10, bg=RGBColor(0x21,0x96,0x53), fg=WHITE, bold=True)

# 범례
tb = s.shapes.add_textbox(Inches(0.3), Inches(6.9), Inches(12), Inches(0.4))
add_text(tb.text_frame,
         "① 사용자 질문 → Bot  ② Coupa 계약 조회  ③ Bedrock 답변  ④ Teams 회신  "
         "⑤ D-60 만료 조회  ⑥ webhook 조회  ⑦ 만료 알림  ⑧ 상태변경 수신  ⑨ 상태 알림",
         9, color=DARK_GRAY)


# ── 슬라이드 3: 시나리오 1 - 질의응답 시퀀스 ─────────────────────
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.5))
add_text(tb.text_frame, "시나리오 1 — Teams 질의응답", 22, bold=True, color=DARK_BLUE)

tb2 = s.shapes.add_textbox(Inches(0.4), Inches(0.65), Inches(12), Inches(0.35))
add_text(tb2.text_frame, "사용자가 Teams에서 봇을 멘션하면 Coupa 계약 데이터를 조회해 Bedrock이 자연어로 답변합니다.", 11, color=DARK_GRAY)

# 참여자 헤더
actors = ["사용자\n(Teams)", "Lambda 1\nBot Handler", "Coupa\nREST API", "Amazon Bedrock\nClaude 3.5"]
colors  = [MID_BLUE, ORANGE, RGBColor(0x21,0x96,0x53), GREEN]
xs = [Inches(0.5), Inches(3.5), Inches(6.8), Inches(10.0)]
for i, (actor, color, x) in enumerate(zip(actors, colors, xs)):
    add_box(s, x, Inches(1.1), Inches(2.5), Inches(0.65), actor, 11, bg=color, fg=WHITE, bold=True)

# 시퀀스 스텝
steps = [
    (0, 1, '@봇 "123번 계약 만료일이 언제야?"'),
    (1, 1, '멘션 태그 제거 / 숫자 ID 추출'),
    (1, 2, 'GET /api/contracts/123?fields=...'),
    (2, 1, '계약 데이터 반환 (name, status, stop_date...)'),
    (1, 3, '질문 + 계약 데이터 컨텍스트 전달'),
    (3, 1, '자연어 답변 생성'),
    (1, 0, '"123번 계약은 2026-06-30 만료입니다."'),
]
y_start = Inches(2.0)
step_h  = Inches(0.62)
for i, (src, dst, label) in enumerate(steps):
    y = y_start + i * step_h
    x_src = xs[src] + Inches(1.25)
    x_dst = xs[dst] + Inches(1.25)
    # 화살표 선
    line = s.shapes.add_connector(1, x_src, y + Inches(0.25), x_dst, y + Inches(0.25))
    line.line.color.rgb = MID_BLUE; line.line.width = Pt(1.5)
    # 라벨
    lx = min(x_src, x_dst) + Inches(0.1)
    lw = abs(x_src - x_dst) - Inches(0.2)
    tb = s.shapes.add_textbox(lx, y, lw, Inches(0.28))
    add_text(tb.text_frame, label, 9, color=DARK_GRAY)


# ── 슬라이드 4: 시나리오 2 - D-60 만료 알림 ─────────────────────
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.5))
add_text(tb.text_frame, "시나리오 2 — 계약 만료 D-60 알림", 22, bold=True, color=DARK_BLUE)

tb2 = s.shapes.add_textbox(Inches(0.4), Inches(0.65), Inches(12), Inches(0.35))
add_text(tb2.text_frame, "매일 오전 9시(KST) EventBridge가 Lambda를 트리거해 만료 임박 계약을 담당자 채널로 알립니다.", 11, color=DARK_GRAY)

actors2 = ["EventBridge\n(매일 9시)", "Lambda 2\nScheduler", "Coupa\nREST API", "DynamoDB\nChannelMapping", "Teams\n채널"]
colors2  = [GREEN, ORANGE, RGBColor(0x21,0x96,0x53), GREEN, MID_BLUE]
xs2 = [Inches(0.3), Inches(2.9), Inches(5.5), Inches(8.1), Inches(10.7)]
for actor, color, x in zip(actors2, colors2, xs2):
    add_box(s, x, Inches(1.1), Inches(2.2), Inches(0.65), actor, 10, bg=color, fg=WHITE, bold=True)

steps2 = [
    (0, 1, '매일 트리거'),
    (1, 2, 'GET /api/contracts?status=active&stop_date[lt_or_eq]=D+60'),
    (2, 1, '만료 임박 계약 목록 반환'),
    (1, 3, 'get_item(PK = 담당자 이메일)'),
    (3, 1, 'webhook URL 반환 (없으면 default fallback)'),
    (1, 4, '⚠️ 만료 60일 전 Adaptive Card 전송'),
]
y_start = Inches(2.0)
for i, (src, dst, label) in enumerate(steps2):
    y = y_start + i * step_h
    x_src = xs2[src] + Inches(1.1)
    x_dst = xs2[dst] + Inches(1.1)
    line = s.shapes.add_connector(1, x_src, y + Inches(0.25), x_dst, y + Inches(0.25))
    line.line.color.rgb = ORANGE; line.line.width = Pt(1.5)
    lx = min(x_src, x_dst) + Inches(0.1)
    lw = abs(x_src - x_dst) - Inches(0.2)
    tb = s.shapes.add_textbox(lx, y, lw, Inches(0.28))
    add_text(tb.text_frame, label, 9, color=DARK_GRAY)


# ── 슬라이드 5: 시나리오 3 - 상태변경 알림 ───────────────────────
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.5))
add_text(tb.text_frame, "시나리오 3 — 계약 상태변경 알림", 22, bold=True, color=DARK_BLUE)

tb2 = s.shapes.add_textbox(Inches(0.4), Inches(0.65), Inches(12), Inches(0.35))
add_text(tb2.text_frame, "Coupa에서 계약 승인/거절/반려/검토 시 담당자 Teams 채널로 즉시 알림을 전송합니다.", 11, color=DARK_GRAY)

actors3 = ["Coupa\nWebhook", "Lambda 3\nWebhook Handler", "DynamoDB\nChannelMapping", "Teams\n채널"]
colors3  = [RGBColor(0x21,0x96,0x53), ORANGE, GREEN, MID_BLUE]
xs3 = [Inches(0.5), Inches(3.8), Inches(7.1), Inches(10.4)]
for actor, color, x in zip(actors3, colors3, xs3):
    add_box(s, x, Inches(1.1), Inches(2.8), Inches(0.65), actor, 11, bg=color, fg=WHITE, bold=True)

steps3 = [
    (0, 1, 'POST /coupa/webhook  (계약 객체 전송)'),
    (1, 1, 'status 필터: approved / rejected / cancelled / pending_approval'),
    (1, 2, 'get_item(PK = 담당자 이메일)'),
    (2, 1, 'webhook URL 반환 (없으면 default fallback)'),
    (1, 3, '✅❌🚫🕐 상태변경 Adaptive Card 전송'),
]
y_start = Inches(2.1)
step_h2 = Inches(0.72)
for i, (src, dst, label) in enumerate(steps3):
    y = y_start + i * step_h2
    x_src = xs3[src] + Inches(1.4)
    x_dst = xs3[dst] + Inches(1.4)
    line = s.shapes.add_connector(1, x_src, y + Inches(0.25), x_dst, y + Inches(0.25))
    line.line.color.rgb = GREEN; line.line.width = Pt(1.5)
    lx = min(x_src, x_dst) + Inches(0.1)
    lw = abs(x_src - x_dst) - Inches(0.2)
    tb = s.shapes.add_textbox(lx, y, lw, Inches(0.28))
    add_text(tb.text_frame, label, 9, color=DARK_GRAY)


# ── 슬라이드 6: DynamoDB 채널 매핑 + Teams 알림 예시 ─────────────
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

tb = s.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12), Inches(0.5))
add_text(tb.text_frame, "DynamoDB 채널 매핑 & Teams 알림 예시", 22, bold=True, color=DARK_BLUE)

# 테이블 헤더
headers = ["email (PK)", "webhook_url", "channel_name"]
col_w = [Inches(3.2), Inches(5.5), Inches(2.0)]
col_x = [Inches(0.4), Inches(3.7), Inches(9.3)]
for hdr, w, x in zip(headers, col_w, col_x):
    add_box(s, x, Inches(0.85), w, Inches(0.45), hdr, 11, bg=DARK_BLUE, fg=WHITE, bold=True, radius=False)

rows = [
    ("hong@company.com",  "https://xxx.webhook.office.com/...", "구매팀"),
    ("kim@company.com",   "https://yyy.webhook.office.com/...", "IT팀"),
    ("lee@company.com",   "https://zzz.webhook.office.com/...", "재무팀"),
]
row_colors = [WHITE, LIGHT_BLUE]
for r, (email, url, ch) in enumerate(rows):
    bg = row_colors[r % 2]
    for val, w, x in zip([email, url, ch], col_w, col_x):
        add_box(s, x, Inches(1.35 + r * 0.45), w, Inches(0.42), val, 10, bg=bg, fg=DARK_GRAY, radius=False)

# 알림 예시 카드 2개
def card(slide, x, y, w, title, title_color, rows_data):
    add_box(slide, x, y, w, Inches(0.45), title, 11, bg=title_color, fg=WHITE, bold=True)
    for i, (k, v) in enumerate(rows_data):
        add_box(slide, x, y + Inches(0.45 + i * 0.38), w / 2, Inches(0.36), k, 9, bg=LIGHT_BLUE, fg=DARK_BLUE, radius=False)
        add_box(slide, x + w / 2, y + Inches(0.45 + i * 0.38), w / 2, Inches(0.36), v, 9, bg=WHITE, fg=DARK_GRAY, radius=False)

card(s, Inches(0.4), Inches(2.55), Inches(5.8),
     "⚠️ 계약 만료 60일 전 알림", ORANGE,
     [("계약명", "SW 라이선스 계약"), ("계약번호", "123"),
      ("만료일", "2026-06-30"), ("담당자", "홍길동"), ("공급업체", "ABC Corp")])

card(s, Inches(6.8), Inches(2.55), Inches(6.1),
     "✅ 계약 승인", GREEN,
     [("계약명", "SW 라이선스 계약"), ("계약번호", "123"),
      ("만료일", "2026-06-30"), ("담당자", "홍길동"),
      ("담당자 이메일", "hong@company.com"), ("공급업체", "ABC Corp")])


out = "chatbot/coupa_teams_chatbot_architecture.pptx"
prs.save(out)
print(f"저장 완료: {out}")
