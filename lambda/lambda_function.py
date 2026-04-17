"""
Coupa AI Facilitator - Lambda Handler
- Microsoft Teams Bot Framework 웹훅 처리
- 기존 Chat UI (POST /chat) 병행 지원
- S3 PDF/PPT 파싱 (콜드스타트 캐싱)
- Bedrock Claude 3 Sonnet
"""
import json, os, io, hmac, hashlib, base64
import boto3

s3      = boto3.client("s3")
bedrock = boto3.client("bedrock-runtime", region_name=os.environ.get("AWS_REGION", "us-east-1"))

MODEL_ID    = "anthropic.claude-3-sonnet-20240229-v1:0"
DOCS_BUCKET = os.environ["DOCS_BUCKET"]
DOCS_PREFIX = os.environ.get("DOCS_PREFIX", "manuals/")
BOT_PASSWORD = os.environ.get("BOT_PASSWORD", "")   # Azure Bot 앱 비밀번호 (서명 검증용)
MAX_TOKENS  = 1024

_doc_cache: str | None = None


# ── 문서 로딩 (S3 캐싱) ────────────────────────────────────────────────────

def get_docs() -> str:
    global _doc_cache
    if _doc_cache is not None:
        return _doc_cache
    pages = s3.list_objects_v2(Bucket=DOCS_BUCKET, Prefix=DOCS_PREFIX)
    parts = []
    for obj in pages.get("Contents", []):
        key = obj["Key"]
        if key.endswith("/"):
            continue
        raw  = s3.get_object(Bucket=DOCS_BUCKET, Key=key)["Body"].read()
        name = key.rsplit("/", 1)[-1]
        try:
            parts.append(_parse(raw, name))
        except Exception as e:
            parts.append(f"[{name} 파싱 실패: {e}]")
    _doc_cache = "\n\n".join(parts) or "등록된 문서가 없습니다."
    return _doc_cache


def _parse(data: bytes, name: str) -> str:
    low = name.lower()
    if low.endswith(".pdf"):
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        text   = "\n".join(p.extract_text() or "" for p in reader.pages)
        return f"=== {name} ===\n{text}"
    if low.endswith((".pptx", ".ppt")):
        from pptx import Presentation
        prs    = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            txts = [s.text.strip() for s in slide.shapes if s.has_text_frame and s.text.strip()]
            if txts:
                slides.append(f"[슬라이드 {i}] " + " | ".join(txts))
        return f"=== {name} ===\n" + "\n".join(slides)
    raise ValueError("지원하지 않는 형식")


# ── Bedrock 호출 ───────────────────────────────────────────────────────────

SYSTEM = """당신은 유니클로 구매/계약 프로세스 안내 챗봇입니다.
반드시 아래 [매뉴얼 문서] 내용만 근거로 답변하세요.
문서에 없는 내용은 "담당 부서에 확인이 필요합니다"라고 답하세요.
답변은 한국어, 마크다운 형식으로 작성하세요.

[매뉴얼 문서]
{docs}"""

def ask_claude(messages: list) -> str:
    result = bedrock.invoke_model(
        modelId=MODEL_ID,
        body=json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": MAX_TOKENS,
            "system": SYSTEM.format(docs=get_docs()),
            "messages": messages,
        }),
    )
    return json.loads(result["body"].read())["content"][0]["text"]


# ── 핸들러 ─────────────────────────────────────────────────────────────────

def lambda_handler(event, context):
    path   = event.get("path", "/")
    method = event.get("httpMethod", "POST")

    # CORS preflight
    if method == "OPTIONS":
        return _resp(200, "")

    body = json.loads(event.get("body") or "{}")

    # ── Teams Bot Framework 웹훅 (/teams)
    if path == "/teams":
        return handle_teams(body)

    # ── 기존 Chat UI (/chat)
    messages = body.get("messages", [])
    if not messages:
        return _resp(400, json.dumps({"error": "messages required"}))
    try:
        reply = ask_claude(messages)
        return _resp(200, json.dumps({"reply": reply}, ensure_ascii=False))
    except Exception as e:
        return _resp(500, json.dumps({"error": str(e)}))


def handle_teams(activity: dict) -> dict:
    """
    Teams Activity 처리
    - type: message  → Claude 답변 후 reply activity 반환
    - type: 기타     → 200 OK (conversationUpdate 등 무시)
    """
    if activity.get("type") != "message":
        return _resp(200, "")

    user_text = activity.get("text", "").strip()
    if not user_text:
        return _resp(200, "")

    try:
        reply_text = ask_claude([{"role": "user", "content": user_text}])
    except Exception as e:
        reply_text = f"오류가 발생했습니다: {e}"

    # Teams reply activity 구성
    reply_activity = {
        "type": "message",
        "text": reply_text,
        "textFormat": "markdown",
        "replyToId": activity.get("id"),
    }
    return _resp(200, json.dumps(reply_activity, ensure_ascii=False))


def _resp(code, body):
    return {
        "statusCode": code,
        "headers": {
            "Content-Type": "application/json",
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Headers": "Content-Type,Authorization",
            "Access-Control-Allow-Methods": "POST,OPTIONS",
        },
        "body": body,
    }
